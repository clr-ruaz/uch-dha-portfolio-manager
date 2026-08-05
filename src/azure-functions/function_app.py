import azure.functions as func
import base64
import csv
import html
import io
import json
import logging
import markdown2
import os
import platform
import pypandoc
import PyPDF2
import re
import tempfile
import zipfile
from bs4 import BeautifulSoup
from docx import Document
from docx.shared import Pt
from docx.oxml.ns import qn
from docx.oxml import OxmlElement
from io import BytesIO
from openpyxl import Workbook
from openpyxl.styles import Font
from openpyxl.utils import get_column_letter
from openpyxl.worksheet.table import Table, TableStyleInfo
from pathlib import Path, PurePosixPath
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import MSO_VERTICAL_ANCHOR, PP_PARAGRAPH_ALIGNMENT
from pptx.util import Emu, Pt

app = func.FunctionApp(http_auth_level=func.AuthLevel.FUNCTION)

MAX_REGEX_PATTERN_LENGTH = 1024
MAX_REGEX_TEXT_LENGTH = 20000

@app.route(route="DqsExtractZip")
def DqsExtractZip(req: func.HttpRequest) -> func.HttpResponse:

    try:
        body = req.get_json()

        if "fileBase64" not in body:
            return func.HttpResponse(
                json.dumps({"error": "Missing fileBase64"}),
                status_code=400,
                mimetype="application/json"
            )

        try:
            zip_bytes = base64.b64decode(body["fileBase64"], validate=True)
        except Exception:
            return func.HttpResponse(
                json.dumps({"error": "Invalid Base64 supplied in 'fileBase64'."}),
                status_code=400,
                mimetype="application/json"
            )

        files = []

        with zipfile.ZipFile(io.BytesIO(zip_bytes), "r") as z:

            for info in z.infolist():

                # Ignore directories
                if info.is_dir():
                    continue

                entry_name = info.filename.replace("\\", "/").lstrip("/")
                entry_path = PurePosixPath(entry_name)
                if any(part == ".." for part in entry_path.parts):
                    continue

                files.append({
                    "path": entry_path.as_posix(),
                    "size": info.file_size,
                    "contentBase64": base64.b64encode(
                        z.read(info.filename)
                    ).decode("utf-8")
                })

        return func.HttpResponse(
            json.dumps({"files": files}),
            mimetype="application/json",
            status_code=200
        )

    except zipfile.BadZipFile:
        return func.HttpResponse(
            json.dumps({"error": "The supplied file is not a valid ZIP archive."}),
            mimetype="application/json",
            status_code=400
        )

    except Exception as ex:
        return func.HttpResponse(
            json.dumps({"error": str(ex)}),
            mimetype="application/json",
            status_code=500
        )


@app.route(route="DqsConvertCsvToXlsx")
def DqsConvertCsvToXlsx(req: func.HttpRequest) -> func.HttpResponse:

    logging.info("Starting DqsConvertCsvToXlsx.")

    try:
        body = req.get_json()

        file_base64 = body.get("fileBase64")
        if not file_base64:
            return func.HttpResponse(
                json.dumps({"success": False, "error": "Missing required property 'fileBase64'."}),
                status_code=400,
                mimetype="application/json"
            )

        file_name = body.get("fileName", "Converted.csv")
        first_row_is_header = bool(body.get("firstRowIsHeader", False))
        format_as_table = bool(body.get("formatAsTable", False))

        try:
            csv_bytes = base64.b64decode(file_base64, validate=True)
        except Exception:
            return func.HttpResponse(
                json.dumps({"success": False, "error": "Invalid Base64 supplied in 'fileBase64'."}),
                status_code=400,
                mimetype="application/json"
            )

        # Decode UTF-8 (with BOM), then fall back to Windows-1252.
        try:
            csv_text = csv_bytes.decode("utf-8-sig")
        except UnicodeDecodeError:
            try:
                csv_text = csv_bytes.decode("cp1252")
            except UnicodeDecodeError:
                return func.HttpResponse(
                    json.dumps({"success": False, "error": "CSV encoding is not supported. Expected UTF-8 or Windows-1252."}),
                    status_code=400,
                    mimetype="application/json"
                )

        # Empty file?
        if not csv_text.strip():
            logging.warning("CSV file '%s' is empty.", file_name)
            return func.HttpResponse(
                json.dumps({"success": False, "error": "The CSV file is empty."}),
                status_code=400,
                mimetype="application/json"
            )

        workbook = Workbook()
        ws = workbook.active
        ws.title = "Data"

        rows_written = 0

        for row in csv.reader(io.StringIO(csv_text)):
            if not any(str(c).strip() for c in row):
                continue
            ws.append(row)
            rows_written += 1

        if rows_written == 0:
            workbook.close()
            logging.warning("CSV file '%s' contains only blank rows.", file_name)
            return func.HttpResponse(
                json.dumps({"success": False,
                            "error": "The CSV file contains only blank rows."}),
                status_code=400,
                mimetype="application/json"
            )

        if first_row_is_header:
            seen = {}
            for c in range(1, ws.max_column + 1):
                cell = ws.cell(row=1, column=c)
                header = "" if cell.value is None else str(cell.value)
                header = header.replace("\\r", " ").replace("\\n", " ")
                header = " ".join(header.split())

                if not header:
                    header = f"Column{c}"

                if header in seen:
                    seen[header] += 1
                    header = f"{header}_{seen[header]}"
                else:
                    seen[header] = 1

                cell.value = header
                cell.font = Font(bold=True)

            ws.freeze_panes = "A2"

        for col in ws.columns:
            width = 10
            for cell in col:
                if cell.value is not None:
                    width = max(width, min(len(str(cell.value)) + 2, 60))
            ws.column_dimensions[get_column_letter(col[0].column)].width = width

        if first_row_is_header:
            ref = f"A1:{get_column_letter(ws.max_column)}{ws.max_row}"

            if format_as_table and ws.max_row >= 2:
                table = Table(displayName="DataTable", ref=ref)
                table.tableStyleInfo = TableStyleInfo(
                    name="TableStyleMedium2",
                    showFirstColumn=False,
                    showLastColumn=False,
                    showRowStripes=True,
                    showColumnStripes=False
                )
                ws.add_table(table)
            else:
                ws.auto_filter.ref = ref

        out = io.BytesIO()
        workbook.save(out)
        workbook.close()

        return func.HttpResponse(
            json.dumps({
                "success": True,
                "fileName": f"{Path(file_name).stem}.xlsx",
                "fileBase64": base64.b64encode(out.getvalue()).decode("utf-8")
            }),
            status_code=200,
            mimetype="application/json"
        )

    except Exception as ex:
        logging.exception("Unexpected error in DqsConvertCsvToXlsx.")
        return func.HttpResponse(
            json.dumps({"success": False, "error": str(ex)}),
            status_code=500,
            mimetype="application/json"
        )


@app.route(route="DqsSplitPdf")
def DqsSplitPdf(req: func.HttpRequest) -> func.HttpResponse:
    logging.info('Python HTTP trigger function processed a request.')

    try:
        jsonRequest_Body = req.get_json()
        strFile_SourceName = jsonRequest_Body.get('fileName') or 'document.pdf'
        b64File_SourcePdf = jsonRequest_Body.get('fileContent')

        if not b64File_SourcePdf:
            return func.HttpResponse(
                json.dumps({"error": "Missing 'fileContent' in request body."}),
                status_code=400,
                mimetype="application/json"
            )

        try:
            binFile_SourcePdf = base64.b64decode(b64File_SourcePdf.encode('ascii'), validate=True)
        except Exception:
            return func.HttpResponse(
                json.dumps({"error": "Invalid Base64 supplied in 'fileContent'."}),
                status_code=400,
                mimetype="application/json"
            )

        safe_source_name = Path(Path(strFile_SourceName).name).stem or 'document'
        strFile_SourceName_NoExt = f"{safe_source_name}_"

        with io.BytesIO(binFile_SourcePdf) as objFile_Input:
            objPdf_SourcePdf = PyPDF2.PdfReader(objFile_Input)
            intFile_SourcePdf_Pages = len(objPdf_SourcePdf.pages)
            logging.info(f"Detected {intFile_SourcePdf_Pages} pages in {strFile_SourceName}.")

            aryFile_PdfOutput = []
            for i in range(intFile_SourcePdf_Pages):
                logging.info(f"Processing PDF page {i}.")
                objPdf_OutputPdf = PyPDF2.PdfWriter()
                objPdf_OutputPdf.add_page(objPdf_SourcePdf.pages[i])

                strFile_OutputName = f"{strFile_SourceName_NoExt}{i + 1}.pdf"
                output_buffer = io.BytesIO()
                objPdf_OutputPdf.write(output_buffer)
                b64File_OutputPdf = base64.b64encode(output_buffer.getvalue()).decode()

                aryFile_PdfOutput.append({
                    'fileName': strFile_OutputName,
                    'fileContent': b64File_OutputPdf,
                    'index': i
                })

        jsonResponse_Body = {
            'Splits': aryFile_PdfOutput
        }

        return func.HttpResponse(body=json.dumps(jsonResponse_Body), mimetype="application/json")

    except Exception as ex:
        logging.exception("Error processing DqsSplitPdf request.")
        return func.HttpResponse(
            json.dumps({"error": str(ex)}),
            status_code=500,
            mimetype="application/json"
        )



@app.route(route="DqsFillPpt", auth_level=func.AuthLevel.FUNCTION)
def DqsFillPpt(req: func.HttpRequest) -> func.HttpResponse:
    logging.info("PowerPoint generation function triggered.")

    try:
        req_body = req.get_json()
        b64_file = req_body.get("fileBase64")
        slides_data = req_body.get("slides_data")

        if not b64_file:
            return func.HttpResponse("Missing 'fileBase64'.", status_code=400)
        if not slides_data:
            return func.HttpResponse("Missing 'slides_data' in request body.", status_code=400)

        try:
            pptx_binary = BytesIO(base64.b64decode(b64_file.encode("utf-8"), validate=True))
        except Exception as decode_error:
            return func.HttpResponse(f"Error decoding base64: {str(decode_error)}", status_code=400)

        prs = Presentation(pptx_binary)

        def get_layout_by_name(prs, layout_name):
            for master in prs.slide_masters:
                for layout in master.slide_layouts:
                    if layout.name == layout_name:
                        idx_name_map = {
                            ph.placeholder_format.idx: ph.name for ph in layout.placeholders
                        }
                        return layout, idx_name_map
            raise ValueError(f"Layout '{layout_name}' not found in any slide master.")

        for slide_data in slides_data:
            layout_name = slide_data["slide_layout_name"]
            placeholder_texts = slide_data["placeholder_names"]
            slide_action = slide_data.get("slide_action", {})
            action_type = slide_action.get("type", "append")
            index = slide_action.get("index", None)

            slide_count = len(prs.slides)

            # --- Fill existing slide ---
            if action_type == "fill":
                if not isinstance(index, int) or not (0 <= index < slide_count):
                    return func.HttpResponse("Invalid 'index' for 'fill' action.", status_code=400)
                slide = prs.slides[index]
                idx_name_map = get_idx_name_map_from_layout(slide.slide_layout)
                for shape in slide.shapes:
                    if not shape.is_placeholder:
                        continue
                    idx = shape.placeholder_format.idx
                    if idx in idx_name_map:
                        original_name = idx_name_map[idx]
                        shape.name = original_name[2:] if original_name.startswith("T_") else original_name
                for shape in slide.shapes:
                    if not shape.is_placeholder:
                        continue
                    ph_name = shape.name
                    if ph_name in placeholder_texts:
                        apply_markdown_to_shape(shape, placeholder_texts[ph_name])
                continue

            # --- Insert from last position ---
            if action_type == "insertLast":
                if not isinstance(index, int) or not (0 <= index <= slide_count):
                    return func.HttpResponse("Invalid 'index' for 'insertLast' action.", status_code=400)
                target_index = slide_count - index
                if target_index < 0:
                    target_index = 0
                slide_layout, idx_name_map = get_layout_by_name(prs, layout_name)
                if 0 <= target_index < slide_count:
                    slide = insert_slide(prs, slide_layout, target_index)
                else:
                    slide = prs.slides.add_slide(slide_layout)
                for shape in slide.shapes:
                    if not shape.is_placeholder:
                        continue
                    idx = shape.placeholder_format.idx
                    if idx in idx_name_map:
                        original_name = idx_name_map[idx]
                        shape.name = original_name[2:] if original_name.startswith("T_") else original_name
                for shape in slide.shapes:
                    if not shape.is_placeholder:
                        continue
                    ph_name = shape.name
                    if ph_name in placeholder_texts:
                        apply_markdown_to_shape(shape, placeholder_texts[ph_name])
                continue

            # --- Normal insert or append ---
            slide_layout, idx_name_map = get_layout_by_name(prs, layout_name)
            if action_type == "insert" and isinstance(index, int) and 0 <= index < slide_count:
                slide = insert_slide(prs, slide_layout, index)
            else:
                slide = prs.slides.add_slide(slide_layout)
            for shape in slide.shapes:
                if not shape.is_placeholder:
                    continue
                idx = shape.placeholder_format.idx
                if idx in idx_name_map:
                    original_name = idx_name_map[idx]
                    shape.name = original_name[2:] if original_name.startswith("T_") else original_name
            for shape in slide.shapes:
                if not shape.is_placeholder:
                    continue
                ph_name = shape.name
                if ph_name in placeholder_texts:
                    apply_markdown_to_shape(shape, placeholder_texts[ph_name])

        output_stream = BytesIO()
        prs.save(output_stream)
        output_stream.seek(0)
        b64_output = base64.b64encode(output_stream.read()).decode("utf-8")

        return func.HttpResponse(
            json.dumps({"fileBase64": b64_output}),
            status_code=200,
            mimetype="application/json"
        )

    except Exception as e:
        logging.error(f"Error in PowerPoint generation: {e}")
        return func.HttpResponse(f"Server error: {str(e)}", status_code=500)

# --- Helper to get the idx→name map from any layout (works for both new and existing slides) ---
def get_idx_name_map_from_layout(layout):
    return {ph.placeholder_format.idx: ph.name for ph in layout.placeholders}
    
# --- Helper to insert slide at specific index ---
def insert_slide(prs, slide_layout, index):
    # Add slide normally (append)
    new_slide = prs.slides.add_slide(slide_layout)

    # Get slide list (sldIdLst) element
    sldIdLst = prs.slides._sldIdLst  
    sldId_elems = list(sldIdLst)

    # Move the last slide id element (new slide) to desired index
    new_slide_id = sldId_elems[-1]
    sldIdLst.remove(new_slide_id)
    sldIdLst.insert(index, new_slide_id)

    return new_slide

# --- Your markdown apply helper ---
def apply_markdown_to_shape(shape, markdown_text):
    if not shape.has_text_frame:
        return

    markdown_text = markdown_text.strip()
    html = markdown2.markdown(markdown_text, extras=["tables"])
    soup = BeautifulSoup(html, "html.parser")
    text_frame = shape.text_frame
    text_frame.clear()

    for paragraph in text_frame.paragraphs:
        text_frame._element.remove(paragraph._element)

    contains_arabic = any("\u0600" <= c <= "\u06FF" for c in markdown_text)

    for element in soup.contents:
        if element.name in ["h1", "h2", "h3"]:
            p = text_frame.add_paragraph()
            p.clear()
            run = p.add_run()
            run.text = element.get_text()
            run.font.bold = True
            if contains_arabic:
                p.alignment = PP_PARAGRAPH_ALIGNMENT.RIGHT

        elif element.name == "ul":
            for li in element.find_all("li"):
                p = text_frame.add_paragraph()
                p.text = li.get_text()
                p.level = 0
                if contains_arabic:
                    p.alignment = PP_PARAGRAPH_ALIGNMENT.RIGHT

        elif element.name == "ol":
            for i, li in enumerate(element.find_all("li"), start=1):
                p = text_frame.add_paragraph()
                p.text = f"{i}. {li.get_text()}"
                p.level = 0
                if contains_arabic:
                    p.alignment = PP_PARAGRAPH_ALIGNMENT.RIGHT

        elif element.name == "table":
            add_table_to_slide(shape, element)

        elif element.name in ["p", "div"]:
            p = text_frame.add_paragraph()
            p.clear()
            for part in element.descendants:
                if isinstance(part, str):
                    run = p.add_run()
                    run.text = part
                elif part.name == "strong":
                    run = p.add_run()
                    run.text = part.get_text()
                    run.font.bold = True
                elif part.name == "em":
                    run = p.add_run()
                    run.text = part.get_text()
                    run.font.italic = True
            if contains_arabic:
                p.alignment = PP_PARAGRAPH_ALIGNMENT.RIGHT
        else:
            continue

def add_table_to_slide(shape, table_soup):
    data = []
    is_header_present = False

    thead = table_soup.find('thead')
    if thead:
        is_header_present = True
        header_data = [th.get_text(strip=True) for th in thead.find_all('th')]
        data.append(header_data)

    tbody = table_soup.find('tbody')
    if tbody:
        for tr in tbody.find_all('tr'):
            row_data = [td.get_text(strip=True) for td in tr.find_all('td')]
            data.append(row_data)

    if not data:
        logging.warning("Markdown table was found but contained no data.")
        return

    num_rows = len(data)
    num_cols = max(len(row) for row in data) if data else 0

    if num_cols == 0:
        return

    slide = shape.part.slide
    left, top, width, height = shape.left, shape.top, shape.width, shape.height

    table_shape = slide.shapes.add_table(num_rows, num_cols, left, top, width, height)
    table = table_shape.table

    col_width = Emu(int(width / num_cols))
    for i in range(num_cols):
        table.columns[i].width = col_width

    for r_idx, row_data in enumerate(data):
        for c_idx, cell_text in enumerate(row_data):
            if c_idx >= num_cols:
                continue

            cell = table.cell(r_idx, c_idx)
            cell.text = cell_text
            cell.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE

            p = cell.text_frame.paragraphs[0]
            run = p.runs[0] if p.runs else p.add_run()
            run.font.size = Pt(14)
            run.font.name = "Bahij Janna"

            if is_header_present and r_idx == 0:
                run.font.bold = True
                p.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
            else:
                contains_arabic = any("\u0600" <= c <= "\u06FF" for c in cell_text)
                if contains_arabic:
                    p.alignment = PP_PARAGRAPH_ALIGNMENT.RIGHT
                else:
                    p.alignment = PP_PARAGRAPH_ALIGNMENT.LEFT

    # Remove placeholder shape after inserting the table
    sp = shape._element
    sp.getparent().remove(sp)
    #logging.info(f"Successfully added a table with {num_rows} rows and {num_cols} columns.")


@app.route(route="DqsStrToDoc", auth_level=func.AuthLevel.FUNCTION)
def DqsStrToDoc(req: func.HttpRequest) -> func.HttpResponse:
    logging.info("DqsStrToDoc function triggered.")

    try:
        req_body = req.get_json()
        markdown_text = req_body.get("md")
        font_name = req_body.get("font", "Calibri")
        font_size_str = req_body.get("size", "11 pt")
        direction_override = req_body.get("direction", None)  # Optional override

        if not markdown_text:
            return func.HttpResponse(
                json.dumps({ "error": "Missing 'md' in request body." }),
                status_code=400,
                mimetype="application/json"
            )

        try:
            font_size = Pt(int(font_size_str.strip().split()[0]))
        except:
            font_size = Pt(11)

        # Determine text direction
        if direction_override:
            rtl_enabled = direction_override.upper() == "RTL"
        else:
            rtl_enabled = contains_arabic(markdown_text)

        # Set Pandoc path
        pandoc_path = set_pandoc_path()
        os.environ["PYPANDOC_PANDOC"] = pandoc_path

        # Convert markdown to docx
        with tempfile.NamedTemporaryFile(delete=False, suffix=".docx") as tmpfile:
            output_path = tmpfile.name
            pypandoc.convert_text(markdown_text, 'docx', format='md', outputfile=output_path)

        # Post-process DOCX
        doc = Document(output_path)

        # Style paragraphs
        for para in doc.paragraphs:
            apply_style_to_paragraph(para, font_name, font_size, rtl=rtl_enabled)

        # Style tables
        for table in doc.tables:
            if rtl_enabled:
                set_table_rtl(table)  # Apply RTL table alignment
            for row in table.rows:
                for cell in row.cells:
                    for para in cell.paragraphs:
                        apply_style_to_paragraph(para, font_name, font_size, rtl=rtl_enabled)

        doc.save(output_path)

        with open(output_path, 'rb') as docx_file:
            docx_bytes = docx_file.read()

        os.remove(output_path)

        return func.HttpResponse(
            json.dumps({ "fileBase64": base64.b64encode(docx_bytes).decode("utf-8") }),
            status_code=200,
            mimetype="application/json"
        )

    except Exception as e:
        logging.error(f"Error: {str(e)}")
        return func.HttpResponse(
            json.dumps({ "error": str(e) }),
            status_code=500,
            mimetype="application/json"
        )

def contains_arabic(text):
    return any('\u0600' <= c <= '\u06FF' or '\u0750' <= c <= '\u08FF' or '\uFB50' <= c <= '\uFEFF' for c in text)

def apply_style_to_paragraph(para, font_name, font_size, rtl=False):
    style_name = para.style.name if para.style else ""

    for run in para.runs:
        run.font.name = font_name
        run._element.rPr.rFonts.set(qn('w:eastAsia'), font_name)

        if style_name in ["Normal", "List Paragraph", ""]:
            run.font.size = font_size

    if rtl:
        pPr = para._element.get_or_add_pPr()
        bidi = OxmlElement('w:bidi')
        pPr.append(bidi)
        rtl_el = OxmlElement('w:rtl')
        rtl_el.set(qn('w:val'), '1')
        pPr.append(rtl_el)

def set_table_rtl(table):
    # Align the table to the right
    tbl_pr = table._element.tblPr
    if tbl_pr is None:
        tbl_pr = OxmlElement('w:tblPr')
        table._element.insert(0, tbl_pr)

    jc = OxmlElement('w:jc')
    jc.set(qn('w:val'), 'right')  # Align table to the right
    tbl_pr.append(jc)

    # Apply RTL to each paragraph inside cells
    for row in table.rows:
        for cell in row.cells:
            for para in cell.paragraphs:
                pPr = para._element.get_or_add_pPr()

                # Set bidi (RTL)
                bidi = OxmlElement('w:bidi')
                bidi.set(qn('w:val'), '1')
                pPr.append(bidi)

                # Optional: set right alignment
                align = OxmlElement('w:jc')
                align.set(qn('w:val'), 'right')
                pPr.append(align)
                
                        
@app.route(route="DqsExtractPpt", auth_level=func.AuthLevel.FUNCTION)
def DqsExtractPpt(req: func.HttpRequest) -> func.HttpResponse:
    logging.info("DqsExtractPpt function triggered.")

    try:
        req_json = req.get_json()

        file_b64 = req_json.get("fileBase64")
        output_format = req_json.get("output_format", "markdown").lower()

        if not file_b64:
            return func.HttpResponse("Missing 'fileBase64' in request body.", status_code=400)

        # Decode base64 PPTX content
        try:
            pptx_bytes = base64.b64decode(file_b64, validate=True)
        except Exception as decode_error:
            return func.HttpResponse(
                f"Error decoding base64: {str(decode_error)}",
                status_code=400,
                mimetype="application/json"
            )
        pptx_stream = io.BytesIO(pptx_bytes)
        prs = Presentation(pptx_stream)

        slides = []
        for idx, slide in enumerate(prs.slides, 1):
            slide_data = {"slide_number": idx, "title": "", "content": []}
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                text_frame = shape.text_frame
                if shape == slide.shapes.title:
                    slide_data["title"] = text_frame.text.strip()
                else:
                    for paragraph in text_frame.paragraphs:
                        slide_data["content"].append({
                            "level": paragraph.level,
                            "text": paragraph.text.strip()
                        })
            slides.append(slide_data)

        if output_format == "markdown":
            result = generate_markdown(slides)
            return func.HttpResponse(result, mimetype="text/markdown")
        elif output_format == "html":
            result = generate_html(slides)
            return func.HttpResponse(result, mimetype="text/html")
        elif output_format == "json":
            result = json.dumps(slides, indent=2, ensure_ascii=False)
            return func.HttpResponse(result, mimetype="application/json")
        else:
            return func.HttpResponse(
                "Invalid 'output_format'. Choose 'Markdown', 'HTML', or 'JSON'.", status_code=400)

    except Exception as e:
        logging.exception("Error processing request in DqsExtractPpt.")
        return func.HttpResponse(f"Internal server error: {str(e)}", status_code=500)


def generate_markdown(slides):
    lines = []
    for slide in slides:
        lines.append(f"# Slide {slide['slide_number']}")
        if slide["title"]:
            lines.append(f"## {slide['title']}")
        for item in slide["content"]:
            indent = "  " * item["level"]
            bullet = "*" if item["level"] > 0 else "-"
            lines.append(f"{indent}{bullet} {item['text']}")
        lines.append("")
    return "\n".join(lines)

def generate_html(slides):
    html_lines = []
    for slide in slides:
        html_lines.append(f"<h1>Slide {slide['slide_number']}</h1>")
        if slide["title"]:
            html_lines.append(f"<h2>{html.escape(slide['title'])}</h2>")
        for item in slide["content"]:
            indent = "&nbsp;" * 4 * item["level"]
            html_lines.append(f"<p>{indent}• {html.escape(item['text'])}</p>")
    return "\n".join(html_lines)

def set_pandoc_path():
    current_dir = os.path.dirname(__file__)
    
    if platform.system() == "Windows":
        # Windows expects .exe and maybe different path, e.g. tools\pandoc.exe
        pandoc_path = os.path.join(current_dir, "tools", "pandoc.exe")
    else:
        # Linux / macOS use the linux binary without extension
        pandoc_path = os.path.join(current_dir, "tools", "pandoc")
    
    # Set environment variable for pypandoc to find pandoc
    os.environ["PYPANDOC_PANDOC"] = pandoc_path
    
    return pandoc_path

@app.route(route="DqsRegEx", auth_level=func.AuthLevel.FUNCTION)
def DqsRegEx(req: func.HttpRequest) -> func.HttpResponse:
    logging.info("Processing DqsRegEx request.")

    try:
        req_body = req.get_json()

        mode = req_body.get("mode")
        pattern = req_body.get("regular_expression")
        flag_str = req_body.get("flag", "")
        original_text = req_body.get("original_text")
        new_text = req_body.get("new_text", "")

        # Map string flags to re flags
        flags = 0
        if flag_str:
            flag_mapping = {
                "IGNORECASE": re.IGNORECASE,
                "MULTILINE": re.MULTILINE,
                "DOTALL": re.DOTALL,
                "UNICODE": re.UNICODE,
                "ASCII": re.ASCII,
                "VERBOSE": re.VERBOSE
            }
            for f in flag_str.upper().split("|"):
                flags |= flag_mapping.get(f.strip(), 0)

        if not all([mode, pattern, original_text]) or mode not in ["match", "replace"]:
            return func.HttpResponse(
                json.dumps({"error": "Invalid input parameters."}),
                status_code=400,
                mimetype="application/json"
            )

        if len(pattern) > MAX_REGEX_PATTERN_LENGTH or len(original_text) > MAX_REGEX_TEXT_LENGTH:
            return func.HttpResponse(
                json.dumps({"error": "Input size exceeds maximum allowed limit."}),
                status_code=400,
                mimetype="application/json"
            )

        try:
            compiled_pattern = re.compile(pattern, flags)
        except re.error as err:
            return func.HttpResponse(
                json.dumps({"error": f"Invalid regular expression: {err}"}),
                status_code=400,
                mimetype="application/json"
            )

        response = {
            "original_text": original_text,
            "updated_text": None,
            "matches": []
        }

        if mode == "match":
            matches = list(compiled_pattern.finditer(original_text))
            response["matches"] = [
                {"match": m.group(0), "start": m.start(), "end": m.end()}
                for m in matches
            ]

        elif mode == "replace":
            matches = list(compiled_pattern.finditer(original_text))
            response["matches"] = [
                {"match": m.group(0), "start": m.start(), "end": m.end()}
                for m in matches
            ]
            updated_text = compiled_pattern.sub(new_text, original_text)
            response["updated_text"] = updated_text

        return func.HttpResponse(
            json.dumps(response, ensure_ascii=False),
            mimetype="application/json"
        )

    except Exception as e:
        logging.error(f"Exception occurred: {str(e)}")
        return func.HttpResponse(
            json.dumps({"error": str(e)}),
            status_code=500,
            mimetype="application/json"
        )


@app.route(route="DqsReplaceTxtInPpt", auth_level=func.AuthLevel.FUNCTION)
def DqsReplaceTxtInPpt(req: func.HttpRequest) -> func.HttpResponse:
    logging.info("Processing DqsReplaceTxtInPpt request.")
    
    """
    Request JSON:
    {
      "fileBase64": "<base64 of .pptx>",          // required
      "replacements": {"old":"new", ...},         // required (str->str)
      "includeMastersLayouts": true,              // optional (default: true)
      "replaceAcrossRuns": false                  // optional (default: false)
    }

    Response JSON:
    {
      "fileBase64": "<base64 of updated .pptx>",
      "bytes": 123456
    }
    """
    try:
        body = req.get_json()
    except Exception:
        return func.HttpResponse("Invalid JSON body.", status_code=400)

    b64_file = body.get("fileBase64")
    raw_repl = body.get("replacements")
    include_ml = bool(body.get("includeMastersLayouts", True))
    replace_across_runs = bool(body.get("replaceAcrossRuns", False))

    if not b64_file:
        return func.HttpResponse("Missing 'fileBase64'.", status_code=400)
    if not isinstance(raw_repl, dict) or not raw_repl:
        return func.HttpResponse("Missing or invalid 'replacements' (expecting a non-empty object).", status_code=400)

    replacements = _normalize_replacements(raw_repl)
    if not replacements:
        return func.HttpResponse("'replacements' normalized to empty after coercion; provide at least one non-empty key.", status_code=400)

    try:
        pptx_bytes = base64.b64decode(b64_file, validate=True)
    except Exception:
        return func.HttpResponse("Invalid base64 in 'fileBase64'.", status_code=400)

    try:
        logging.info(
            "Starting PPTX replace operation | bytes=%d | includeML=%s | acrossRuns=%s",
            len(pptx_bytes), include_ml, replace_across_runs
        )
        updated_bytes = replace_everywhere_in_memory(
            pptx_bytes,
            replacements,
            include_masters_layouts=include_ml,
            replace_across_runs=replace_across_runs
        )
        logging.info("Completed PPTX replace operation | outBytes=%d", len(updated_bytes))
    except Exception as e:
        logging.exception("Replace operation failed.")
        return func.HttpResponse(f"Replace operation failed: {e}", status_code=500)

    resp = {
        "fileBase64": base64.b64encode(updated_bytes).decode("utf-8"),
        "bytes": len(updated_bytes),
    }
    return func.HttpResponse(json.dumps(resp), status_code=200, mimetype="application/json")


# ----------------------
# Helper functions below
# ----------------------

def _normalize_replacements(raw: dict) -> dict:
    """Coerce keys/values to strings; drop empties."""
    norm = {}
    for k, v in raw.items():
        ks = "" if k is None else str(k)
        vs = "" if v is None else str(v)
        if ks:
            norm[ks] = vs
    return norm


def replace_in_text_frame_runsafe(text_frame, replacements: dict):
    """Replace within each run (preserves formatting but misses cross-run matches)."""
    if not text_frame:
        return
    for p in text_frame.paragraphs:
        for r in p.runs:
            t = r.text or ""
            if not t:
                continue
            for old, new in replacements.items():
                if old in t:
                    t = t.replace(old, new)
            r.text = t


def replace_in_text_frame_paragraph(text_frame, replacements: dict):
    """
    Replace across the entire paragraph text (catches cross-run matches)
    but will reset character-level formatting within each paragraph.
    """
    if not text_frame:
        return
    for p in text_frame.paragraphs:
        full = "".join([run.text or "" for run in p.runs]) or ""
        if not full:
            continue
        replaced = full
        for old, new in replacements.items():
            if old in replaced:
                replaced = replaced.replace(old, new)
        p.clear()
        run = p.add_run()
        run.text = replaced


def _replace_text_frame(text_frame, replacements: dict, replace_across_runs: bool):
    if replace_across_runs:
        replace_in_text_frame_paragraph(text_frame, replacements)
    else:
        replace_in_text_frame_runsafe(text_frame, replacements)


def handle_table(table, replacements: dict, replace_across_runs: bool):
    if not table:
        return
    for row in table.rows:
        for cell in row.cells:
            tf = getattr(cell, "text_frame", None)
            if tf is not None:
                _replace_text_frame(tf, replacements, replace_across_runs)


def handle_chart(chart, replacements: dict, replace_across_runs: bool):
    if not chart:
        return

    # Chart title
    try:
        ct = getattr(chart, "chart_title", None)
        if getattr(chart, "has_title", False) and ct and getattr(ct, "has_text_frame", False):
            _replace_text_frame(ct.text_frame, replacements, replace_across_runs)
    except Exception:
        logging.exception("Chart title replacement skipped due to structure.")

    # Category axis title
    try:
        if getattr(chart, "has_category_axis", False):
            ax = chart.category_axis
            at = getattr(ax, "axis_title", None)
            if ax and getattr(ax, "has_title", False) and at and getattr(at, "has_text_frame", False):
                _replace_text_frame(at.text_frame, replacements, replace_across_runs)
    except Exception:
        logging.exception("Category axis title replacement skipped.")

    # Value axis title
    try:
        if getattr(chart, "has_value_axis", False):
            ax = chart.value_axis
            at = getattr(ax, "axis_title", None)
            if ax and getattr(ax, "has_title", False) and at and getattr(at, "has_text_frame", False):
                _replace_text_frame(at.text_frame, replacements, replace_across_runs)
    except Exception:
        logging.exception("Value axis title replacement skipped.")

    # Note: Legend entries/data labels not editable as free text via python-pptx.


def walk_shapes(shapes, replacements: dict, replace_across_runs: bool):
    for shp in shapes:
        try:
            if shp.shape_type == MSO_SHAPE_TYPE.GROUP:
                walk_shapes(shp.shapes, replacements, replace_across_runs)
                continue

            if getattr(shp, "has_text_frame", False):
                _replace_text_frame(shp.text_frame, replacements, replace_across_runs)

            if shp.shape_type == MSO_SHAPE_TYPE.TABLE and getattr(shp, "table", None) is not None:
                handle_table(shp.table, replacements, replace_across_runs)

            if getattr(shp, "has_chart", False) and getattr(shp, "chart", None) is not None:
                handle_chart(shp.chart, replacements, replace_across_runs)

            # SmartArt still not exposed via python-pptx.
        except Exception:
            logging.exception("A shape was skipped due to unexpected structure.")


def replace_everywhere_in_memory(pptx_bytes: bytes, replacements: dict, include_masters_layouts: bool, replace_across_runs: bool) -> bytes:
    prs = Presentation(BytesIO(pptx_bytes))

    for slide in prs.slides:
        walk_shapes(slide.shapes, replacements, replace_across_runs)

    if include_masters_layouts:
        for layout in prs.slide_layouts:
            walk_shapes(layout.shapes, replacements, replace_across_runs)
        for master in prs.slide_masters:
            walk_shapes(master.shapes, replacements, replace_across_runs)

    out = BytesIO()
    prs.save(out)
    return out.getvalue()
